#!/usr/bin/env python3
"""
Adapted document renderer for the G12++ suite.

Unlike the original gen.py, this does NOT read a results spreadsheet — the suite
already holds the data. It takes a JSON job (the Student Summary built from the
locked-grades read-model + uploaded templates + per-cycle settings), fills the
PowerPoint templates with python-pptx, converts to PDF with LibreOffice, and zips
per document type. It prints a JSON result on stdout.

Usage:  python3 doc_gen.py <job.json>

job.json:
{
  "kinds": ["certificate", "report"],
  "templates": { "certificate": "/path/cert.pptx", "report": "/path/report.pptx" },
  "settings": { "cycleName": "...", "testCentre": "...", "examDate": "...", "issueDate": "..." },
  "students": [
    { "resultId": "P0001", "name": "Student 01", "award": "Distinction award",
      "subjects": { "S1": {"level": "...", "stars": "***"}, ... "S5": {...} } }
  ],
  "outDir": "/tmp/...", "workDir": "/tmp/...", "fontDir": "/path/.fonts"
}
"""
import os, re, sys, json, shutil, zipfile, subprocess

BARLOW = "https://raw.githubusercontent.com/google/fonts/main/ofl/barlow/{}.ttf"
# Tokenisation of each template's baked-in / placeholder fields into merge
# tokens. Per-student RESULTID and per-cycle TESTCENTRE/EXAMDATE/ISSUEDATE
# replace the certificate's hard-coded values and the report's dash placeholders.
CERT_CONSTANTS = [
    ("A-A-011108", "{{RESULTID}}"),       # baked-in fixed Result ID
    ("Alsama Shatila 1", "{{TESTCENTRE}}"),
    ("11th March 2026", "{{ISSUEDATE}}"),
]
# Report dash placeholders, longest first so a shorter run can't match inside a
# longer one. Each length is unique in the template.
REPORT_DASHES = [
    ("------------", "{{TESTCENTRE}}"),  # 12
    ("----------", "{{RESULTID}}"),      # 10
    ("---------", "{{ISSUEDATE}}"),      # 9
    ("--------", "{{EXAMDATE}}"),        # 8
]


def log(*a):
    print(*a, file=sys.stderr)


from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def all_runs(shapes):
    for sh in shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from all_runs(sh.shapes)
        elif sh.has_table:
            for row in sh.table.rows:
                for cell in row.cells:
                    for p in cell.text_frame.paragraphs:
                        yield from p.runs
        elif sh.has_text_frame:
            for p in sh.text_frame.paragraphs:
                yield from p.runs


def setup_fonts(font_dir):
    """Fetch Barlow; detect that Georgia Pro Condensed (cert name line) is
    absent. Returns a warnings dict the UI surfaces."""
    os.makedirs(font_dir, exist_ok=True)
    warnings = []
    try:
        have = subprocess.run(["fc-list"], capture_output=True, text=True).stdout.lower()
    except Exception:
        have = ""
    barlow_present = "barlow" in have
    if not barlow_present:
        for face in ("Barlow-Regular", "Barlow-Bold"):
            try:
                subprocess.run(
                    ["curl", "-fsSL", "-o", os.path.join(font_dir, face + ".ttf"), BARLOW.format(face)],
                    check=True, timeout=30,
                )
            except Exception:
                warnings.append(f"Could not fetch {face} — reports may fall back to a default sans.")
        try:
            subprocess.run(["fc-cache", "-f", font_dir], capture_output=True, timeout=30)
            have = subprocess.run(["fc-list"], capture_output=True, text=True).stdout.lower()
        except Exception:
            pass
        barlow_present = "barlow" in have
    georgia_present = "georgia pro condensed" in have
    if not georgia_present:
        warnings.append(
            "Georgia Pro Condensed is not installed (it's a proprietary font used on the "
            "certificate name line). The name line will render in a substitute font. To keep "
            "the exact look, embed fonts in the template: PowerPoint → Save → Embed fonts in the file."
        )
    return {"georgiaPresent": georgia_present, "barlowPresent": barlow_present, "warnings": warnings}


def tokenize_run(text, kind):
    """Normalise a template's baked/placeholder values to merge tokens."""
    if kind == "certificate":
        for old, tok in CERT_CONSTANTS:
            if old in text:
                text = text.replace(old, tok)
    else:
        for dashes, tok in REPORT_DASHES:
            if dashes in text:
                text = text.replace(dashes, tok)
    return text


def fill(template, kind, tokens, dest):
    prs = Presentation(template)
    for s in prs.slides:
        for run in all_runs(s.shapes):
            t = tokenize_run(run.text, kind)
            if "{{" in t:
                for key, val in tokens.items():
                    t = t.replace("{{" + key + "}}", val)
            if t != run.text:
                run.text = t
    prs.save(dest)


def safe(name):
    return re.sub(r"[^\w\-. ]", "", name or "").strip().replace(" ", "_") or "student"


def tokens_for(kind, student, settings):
    base = {
        "NAME": student.get("name", ""),
        "RESULTID": student.get("resultId", ""),
        "TESTCENTRE": settings.get("testCentre", ""),
        "EXAMDATE": settings.get("examDate", ""),
        "ISSUEDATE": settings.get("issueDate", ""),
        "CYCLE": settings.get("cycleName", ""),
    }
    if kind == "certificate":
        base["AWARD"] = student.get("award", "")
    else:
        subjects = student.get("subjects", {})
        for slot in ("S1", "S2", "S3", "S4", "S5"):
            s = subjects.get(slot, {})
            base[f"{slot}_LEVEL"] = s.get("level", "")
            base[f"{slot}_STARS"] = s.get("stars", "")
    return base


def soffice_env():
    """LibreOffice in a sandbox needs the headless 'svp' VCL plugin and a
    writable HOME for its profile. Inherits PATH from the parent."""
    env = os.environ.copy()
    env["SAL_USE_VCLPLUGIN"] = "svp"
    env.setdefault("HOME", "/tmp")
    return env


def soffice_convert(pptxs, work):
    """Batch-convert PPTX → PDF with a private LibreOffice profile."""
    profile = "file://" + os.path.join(work, "lo_profile")
    cmd = ["soffice", f"-env:UserInstallation={profile}", "--headless",
           "--convert-to", "pdf", "--outdir", work] + pptxs
    subprocess.run(cmd, check=True, capture_output=True, timeout=600, env=soffice_env())


def build(kind, template, students, settings, work_root, out_dir):
    suffix, zipname = {
        "certificate": ("Certificate", "certificates.zip"),
        "report": ("Report", "performance_reports.zip"),
    }[kind]
    work = os.path.join(work_root, kind)
    shutil.rmtree(work, ignore_errors=True)
    os.makedirs(work)

    per_student = {}
    pptx_for = {}
    for st in students:
        sid = st.get("resultId") or safe(st.get("name"))
        if not (st.get("name") or "").strip():
            per_student[sid] = {"status": "error", "error": "Name is empty in the student summary."}
            continue
        try:
            dest = os.path.join(work, f"{safe(st['name'])}_{sid}_{suffix}.pptx")
            fill(template, kind, tokens_for(kind, st, settings), dest)
            pptx_for[sid] = dest
            per_student[sid] = {"status": "pending"}
        except Exception as e:  # pragma: no cover - defensive
            per_student[sid] = {"status": "error", "error": f"Fill failed: {e}"}

    pptxs = sorted(pptx_for.values())
    if pptxs:
        soffice_convert(pptxs, work)

    zpath = os.path.join(out_dir, zipname)
    with zipfile.ZipFile(zpath, "w", zipfile.ZIP_DEFLATED) as z:
        for sid, src in pptx_for.items():
            pdf = os.path.splitext(src)[0] + ".pdf"
            if os.path.exists(pdf):
                out_pdf = os.path.join(out_dir, os.path.basename(pdf))
                shutil.copy(pdf, out_pdf)
                z.write(pdf, os.path.basename(pdf))
                per_student[sid] = {"status": "complete", "file": os.path.basename(pdf)}
            else:
                per_student[sid] = {"status": "error", "error": "LibreOffice did not produce a PDF."}

    complete = sum(1 for v in per_student.values() if v["status"] == "complete")
    return {"zip": os.path.basename(zpath), "complete": complete, "total": len(students), "perStudent": per_student}


def main(job_path):
    job = json.load(open(job_path))
    out_dir = job["outDir"]
    work_root = job.get("workDir", os.path.join(out_dir, "work"))
    os.makedirs(out_dir, exist_ok=True)
    os.makedirs(work_root, exist_ok=True)
    font_dir = job.get("fontDir", os.path.expanduser("~/.fonts"))

    fonts = setup_fonts(font_dir)
    students = job["students"]
    settings = job.get("settings", {})
    templates = job["templates"]

    result = {"fonts": fonts, "kinds": {}}
    for kind in job["kinds"]:
        tpl = templates.get(kind)
        if not tpl or not os.path.exists(tpl):
            result["kinds"][kind] = {"error": f"No {kind} template provided."}
            continue
        result["kinds"][kind] = build(kind, tpl, students, settings, work_root, out_dir)

    print(json.dumps(result))


if __name__ == "__main__":
    if len(sys.argv) < 2:
        log("usage: doc_gen.py <job.json>")
        sys.exit(2)
    main(sys.argv[1])
